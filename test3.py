from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_academic_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 Widescreen
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # Blank layout to manually position B&W boxes

    # Constants for minimalist styling
    FONT_TITLE = "Arial"
    FONT_BODY = "Arial"
    COLOR_BLACK = RGBColor(0, 0, 0)

    def add_slide(title_text, bullets):
        slide = prs.slides.add_slide(blank_layout)
        
        # Title Box
        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.33), Inches(1.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = FONT_TITLE
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = COLOR_BLACK
        p.alignment = PP_ALIGN.LEFT

        # Content Box
        contentBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(5.0))
        ctf = contentBox.text_frame
        ctf.word_wrap = True
        
        for i, text in enumerate(bullets):
            p = ctf.add_paragraph() if i > 0 else ctf.paragraphs[0]
            if text.startswith("  * ") or text.startswith("    - "): # Multi-level emulation
                p.text = text.replace("  * ", "").replace("    - ", "")
                p.level = 1 if text.startswith("  * ") else 2
            else:
                p.text = text.replace("* ", "")
                p.level = 0
            
            p.font.name = FONT_BODY
            p.font.size = Pt(22) if p.level == 0 else Pt(18)
            p.font.color.rgb = COLOR_BLACK
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(12)

    # -------------------------------------------------------------------------
    # Slide 1: Title Slide
    # -------------------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    titleBox = slide1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.33), Inches(3.0))
    tf = titleBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "HiPert: A Three-Layer Syntax-Semantics-Aware Defense Against Indirect Prompt Injection in LLM-based Code Contexts"
    p.font.name = FONT_TITLE
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLACK
    
    infoBox = slide1.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(11.33), Inches(2.2))
    itf = infoBox.text_frame
    infos = [
        "Student: Jack (Po Han Cheng)",
        "Advisor: Ying-Dar Lin, Ph.D.",
        "High Speed Network Lab",
        "National Yang Ming Chiao Tung University, Taiwan",
        "June 2026"
    ]
    for i, text in enumerate(infos):
        p = itf.add_paragraph() if i > 0 else itf.paragraphs[0]
        p.text = text
        p.font.name = FONT_BODY
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_BLACK

    # -------------------------------------------------------------------------
    # Slide 2 - Slide 34 Data Insertion
    # -------------------------------------------------------------------------
    slides_data = [
        ("Outline", [
            "* Motivation & Background",
            "* Notation & Problem Formulation",
            "* Related Work (Attacks & Defenses)",
            "* HiPert Detection System Architecture",
            "* Core Issues (Q1 - Q5)",
            "* Experiment Setup & Evaluation Results",
            "* Conclusion & References"
        ]),
        ("Motivation (1/2)", [
            "* Top 1 Risk: OWASP LLM security ranking",
            "* Impact: Malicious code generation (vulnerabilities/backdoors)",
            "* Challenge: Hard to locate in massive external code resources"
        ]),
        ("Motivation (2/2)", [
            "* Attack Vector: Embedded in non-executable areas (Comments/Strings)",
            "* Human Developer View: Meaningless noise / Trivial docstrings",
            "* Code LLM View: Attention steering vector",
            "*(Note: Insert Black & White Python Code Snippet Box Here)*"
        ]),
        ("Background – Indirect Prompt Injection Attack", [
            "* Source: Code files (C++, Python), PDF, Word",
            "* Mechanism: Untrusted data treated as command instructions",
            "* Blindspot: Bypasses traditional static code analysis scanners"
        ]),
        ("Indirect Prompt Injection Attack Example (1/2)", [
            "* User Prompt: Context completion request",
            "* External Code: Hidden block (# SYSTEM CONFIGURATION # MODE: SUDO)",
            "* LLM Response: Dangerous script backdoor (os.system(\"sudo rm -rf /\"))"
        ]),
        ("Indirect Prompt Injection Attack Example (2/2)", [
            "* Structure: Full Prompt = Position Code + Noise Code + Target Code",
            "* Optimization: Automated noise perturbation to bypass baseline text scanners"
        ]),
        ("Background – Tokenization & AST / CST", [
            "* Tokenization: Extract raw comments only; lacks syntax hierarchy",
            "* AST (Abstract Syntax Tree): Identifies variables; discards non-executable comments",
            "* CST (Concrete Syntax Tree - Tree-sitter): 100% accurate byte-level spans (Comments, Strings, Identifiers)"
        ]),
        ("Background – Perplexity (PPL)", [
            "* Formula: PPL(W) = N-th_root( 1 / P(w1, w2, ..., wN) )",
            "* Gibberish Payload: High PPL -> Easy to catch",
            "* Natural Semantic Trigger: Low PPL -> Passes undetected",
            "* Conclusion: PPL alone = Insufficient"
        ]),
        ("Background – Min-K%", [
            "* Formula: Min-K% = (1 / |E|) * Sum( log P(xt | x<t) )",
            "* Advantage: Focuses on lowest-probability tokens; reduces benign dilution",
            "* Limitation: Fixed k is unstable across variable code node lengths",
            "* HiPert Optimization: CST-guided dynamic k + Entropy weighting"
        ]),
        ("Notations (1/2)", [
            "* Pu: User Prompt",
            "* C: External Code Files (C = {row_1, ..., row_m})",
            "* A / T: AST Parser (Tree-sitter) / Tokenizer",
            "* Dstr, Dcom, Dvar: Extracted Strings, Comments, Variables",
            "* n: Target CST Node (n in Dcom U Dstr U Dvar)"
        ]),
        ("Notations (2/2)", [
            "* GS / GA: Semantic Guardrail / Adversarial Guardrail",
            "* sem: Semantic Score (Behavioral Shift Metric)",
            "* mink: Min-K% Anomaly Score",
            "* sem, mink: Risk Threshold Boundaries"
        ]),
        ("Problem Overview", [
            "* Input: Untrusted Context x = (Pu, C)",
            "* Benign Path: Directly processed -> Safe Output (ysafe)",
            "* Malicious Path: Detection Engine -> Sanitization ~C = Sanitize(C, S) -> Safe Code Generation"
        ]),
        ("Problem Statement", [
            "* Objective: Maximize F1-Score; Minimize ASR; Minimize Latency",
            "* Constraints:",
            "  * Localization: S subset of Crisk (Ccomment, Cstring, Cidentifier, Cdecoy)",
            "  * Preservation: Core executable logic unchanged (Cexec(~C) = Cexec(C))",
            "  * Minimality: Smallest modification budget (|S| <= m_bar)"
        ]),
        ("Related Work ---- Prompt Injection Attack", [
            "* Attack Taxonomy:",
            "  * Adversarial (ShadowCode, INSEC): Gradient-optimized / Gibberish noise / Unnatural token distribution -> Handled by Layer 2 (Min-K%)",
            "  * Semantic (XOXO, ITGen): Fluent renaming / Natural surface probability / High stealth -> Handled by Layer 3 (Perturbation)",
            "*(Note: Insert standard B&W matrix comparison here)*"
        ]),
        ("Related Work ---- Prompt Injection Defense", [
            "* UniGuardian: NLP Word Filtering",
            "* DePA: Line Perplexity (Line-level analysis)",
            "* KillBadCode: Code Naturalness (Token-level analysis)",
            "* CodeGarrison: Hybrid Vector / Trained Embeddings",
            "* HiPert (Ours): CST Node-level + Layer 2 (Dynamic Min-K%) + Layer 3 (Node Perturbation)"
        ]),
        ("Solution Overview", [
            "* Stage 1: Input Processing & Target Extraction",
            "  * Tree-sitter parsing -> Positional structure map",
            "  * Target nodes: Dcand = Dcom U Dstr U Dvar",
            "* Stage 2: HiPert Detection System (Hierarchical)",
            "  * Layer 1 (GF): Syntax-Guided Pre-Filter (Regex, Structural, Decoy)",
            "  * Layer 2 (GA): CST-Guided Dynamic Min-K% (Adversarial scoring)",
            "  * Layer 3 (GS): CST-Guided Node Perturbation Analysis (Semantic scoring)"
        ]),
        ("Solution Overview – Technical Details", [
            "* Pre-Filtering: Laplace-smoothed character & non-ASCII ratios over nodes",
            "* Decoy Mitigation (Flashboom Defense): Local program dependency call-graph; prunes isolated, unreferenced sections"
        ]),
        ("Solution: Semantic Guardrail and Adversarial Guardrail", [
            "* Adversarial Guardrail (Layer 2):",
            "  * Target: Adversarial (ShadowCode, INSEC) -> Unnatural noise / Gibberish",
            "  * Method: Dynamic Min-K% Scoring (Token NLL + Spike detection)",
            "* Semantic Guardrail (Layer 3):",
            "  * Target: Semantic (XOXO, ITGen) -> Fluent / Natural renaming",
            "  * Method: Node Perturbation Analysis -> Logits distance d(.,.) after neutralization"
        ]),
        ("Issues", [
            "* Q1: Outperform 3 baselines across 6 attack types?",
            "* Q2: Generalize across different surrogate defense models?",
            "* Q3: Are all three layers non-redundant and necessary?",
            "* Q4: Precision vs. Recall discrimination capability?",
            "* Q5: Performance vs. general NLP-based defense methods?"
        ]),
        ("Experiment Setup: Data Source & Test models", [
            "* Attack Datasets: ShadowCode, INSEC, XOXO, ITGen, Flashboom, CoTDeceptor",
            "* Defense Surrogate Models: Codegen-350M, Qwen3.5-4B, Gemma-4-E4B",
            "* Victim Architectures: CodeGemma, CodeBERT, GraphBERT, Mistral-7B, Claude-3.5, ChatGPT-5.1, Gemini-3.1"
        ]),
        ("Evaluation Metric (1/2)", [
            "* Outcomes: TP (Caught), FP (False Alarm), TN (Clean Passed), FN (Missed)",
            "* Standard Metrics: Precision, Recall, F1-Score, False Positive Rate (FPR)"
        ]),
        ("Evaluation Metric (2/2)", [
            "* Attack Success Rate (ASR): ASR = N_success / N_total",
            "* Code Utility: Lexical Similarity (Edit Distance) & Structural Similarity (AST overlap)"
        ]),
        ("Experiment Results: Q1 Detection Performance (1/2)", [
            "* Rank #1: Best average F1-score (0.80) across 6 attacks",
            "* ShadowCode: +33.8% F1 improvement (Dynamic Min-K%)",
            "* XOXO: -45% False Positive Rate reduction",
            "* ITGen: +24% Precision increase",
            "*(Note: Insert Black & White 4-Panel Evaluation Chart Here)*"
        ]),
        ("Experiment Results: Q1 Performance (2/2)", [
            "* AUROC: 0.820 threshold-agnostic performance",
            "* Low-FPR Boundary: Intercepts 52% of attacks at < 5% FPR",
            "*(Note: Insert Black & White ROC Curve Chart Here)*"
        ]),
        ("Experiment Results: Q2 Cross-Model Generalization", [
            "* Surrogate Stability: F1 remains 0.737 - 0.807 across 5 surrogate models",
            "* High Specificity: Precision > 0.84 in 4 out of 5 models",
            "*(Note: Insert Black & White Heatmap Matrix Here)*"
        ]),
        ("Experiment Results: Q3 Ablation Study (1/2)", [
            "* Layer Contributions: L1 (Structural Check) -> L2 (Statistical Noise) -> L3 (Semantic Shift)",
            "* F1 Cumulative Progression: 0.19 (L1) -> 0.39 (L1+L2) -> 0.82 (Final Full Pipeline)",
            "*(Note: Insert Black & White Ablation Progression Bars Here)*"
        ]),
        ("Experiment Results: Q3 Ablation Study (2/2)", [
            "* Latency vs. Performance Trade-off:",
            "  * Layer 3 Standalone: Strongest individual F1 (0.748) but high latency (164.20 ms)",
            "  * HiPert Pipeline Combination: L1 & L2 act as early-exit filters; drops average latency to 104.64 ms, boosts F1 to 0.82"
        ]),
        ("Experiment Results: Adaptive Attack Robustness", [
            "* Strategies Tested: Decoy Injection (-8% F1), Copy Trigger (-20% F1), Contextual Attack (-16% F1)",
            "* Lower Bound: HiPert minimum F1 (0.62) significantly outperforms single baselines (KillBadCode: 0.288)",
            "*(Note: Insert Black & White Robustness Bars Chart Here)*"
        ]),
        ("References (1/3)", [
            "* [1] Ghannoum et al., 'Poisoned Source Code Detection in Code Models,' 2025.",
            "* [2] Tsai et al., 'Beyond Natural Language Perplexity: Detecting Dead Code Poisoning,' 2025.",
            "* [3] Li et al., 'Blinding LLM-Based Code Auditors with Flashboom Attacks,' 2025.",
            "* [4] Huang et al., 'Iterative Generation of Adversarial Example for Deep Code Models,' 2025."
        ]),
        ("References (2/3)", [
            "* [5] Storek et al., 'XOXO: Cross-Origin Context Poisoning Attacks,' 2025.",
            "* [6] Jenko et al., 'Black-Box Adversarial Attacks on LLM-Based Code Completion,' 2025.",
            "* [7] Sun et al., 'Kill Code Poisoning: Lightweight Method Based on Code Naturalness,' 2025."
        ]),
        ("References (3/3)", [
            "* [8] Liu et al., 'Prompt Injection Attacks on Agentic AI Coding Editors,' 2025.",
            "* [9] Wahed et al., 'MOCHA: Are Code Language Models Robust Against Malicious Coding Prompts?' 2025.",
            "* [10] Yang et al., 'ShadowCode: External Prompt Injection Attack against Code LLMs,' 2025."
        ]),
        ("Conclusion & Thank You", [
            "* Core Contribution: CST parsing + behavioral perturbation = Optimal code defense",
            "* ASR Reduction: Dropped downstream attack success from 64%-75% to 10.5%-21.5%",
            "* Generalizability: Robust against black-box victims, adaptive strategies, and tokenizer mismatches",
            "Thank you. Q&A",
            "High Speed Network Lab, NYCU, Taiwan"
        ]),
        ("Appendix: Taxonomy Categorization", [
            "* 1. Semantic Attacks (XOXO, ITGen, Flashboom, CoTDeceptor):",
            "  * Interfere with instruction-following via fluent, natural-looking transformations.",
            "* 2. Adversarial Attacks (GCG, ShadowCode, INSEC):",
            "  * Generate non-functional token perturbations optimized via model gradient searches."
        ])
    ]

    for title, bullets in slides_data:
        add_slide(title, bullets)

    prs.save("HiPert_Oral_Defense.pptx")
    print("B&W Academic PPT Successfully generated: HiPert_Oral_Defense.pptx")

if __name__ == "__main__":
    create_academic_ppt()